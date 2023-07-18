from pptx import Presentation

from pptx.util import Inches, Pt
import matplotlib.pyplot as plt
import json

with open("sample.json", "r") as jsonfile:
    data = json.load(jsonfile)
    print("Read successful")
#presentation lista beolvasása
pres=data['presentation']
#title dia létrehozása
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
#title kiírás
title.text =str(pres[0]["title"])
subtitle.text =str(pres[0]["content"])
#text dia létrehozás
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

left = top = width = height = Inches(1)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame

tf.text = str(pres[1]["title"])

p = tf.add_paragraph()
p.text = str(pres[1]["content"])

#lista dia létrehozás
bullet_slide_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = str(pres[2]["title"])

tf = body_shape.text_frame
tf.text = str(pres[2]["title"])

list=pres[2]["content"]

for i in range(len(list)):
    p = tf.add_paragraph()
    p.text = list[i]["text"]
    p.level = list[i]["level"]
#picture dia létrehozás
img_path = pres[3]["content"]
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

left = Inches(2.5)
height = Inches(5)
pic = slide.shapes.add_picture(img_path, left, top, height=height)
#plot chart létrehozás

blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)



x = [1,3,5.6,8.9,]
y = [2,4,6.7,9.10]

plt.plot(x, y)

plt.xlabel(pres[4]["configuration"]["x-label"])
plt.ylabel(pres[4]["configuration"]["y-label"])

#A generált chartot elmentjük képként a következő néven "plot.png" (lehet változó is)
plt.show()
plt.title(pres[4]["title"])

img_path = "plot.png"
left = Inches(1.5)
height = Inches(5)
pic = slide.shapes.add_picture(img_path, left, top, height=height)
#elmentés
prs.save('test.pptx')
